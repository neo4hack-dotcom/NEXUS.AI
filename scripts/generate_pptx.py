"""
DOINg.AI — User Guide PowerPoint generator
Fully editable 16:9 deck, brand-consistent with the PDF booklet.
Every text element is a proper Shape → editable in PowerPoint/Keynote.

Run:  python3 scripts/generate_pptx.py
Out:  ./DOINg-AI-User-Guide.pptx
"""
from __future__ import annotations
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ── Brand palette ─────────────────────────────────────────────────────────────
BRAND   = RGBColor(0xFF, 0x3E, 0x00)   # #FF3E00
INK_950 = RGBColor(0x0A, 0x0B, 0x0E)   # near-black
INK_900 = RGBColor(0x15, 0x17, 0x1C)
INK_700 = RGBColor(0x2D, 0x30, 0x38)
INK_500 = RGBColor(0x6E, 0x72, 0x80)
INK_300 = RGBColor(0xC2, 0xC5, 0xCF)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
PAPER   = RGBColor(0xFA, 0xFA, 0xF8)
TEXT    = RGBColor(0x0F, 0x11, 0x15)
MUTED   = RGBColor(0x6E, 0x72, 0x80)
BLUE    = RGBColor(0x3B, 0x82, 0xF6)
GREEN   = RGBColor(0x10, 0xB9, 0x81)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)

# ── Slide dimensions (16:9, 10 × 5.625 in) ───────────────────────────────────
W = Inches(10)
H = Inches(5.625)

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]   # completely blank layout
    return prs.slides.add_slide(layout)

# ── Shape helpers ─────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill: RGBColor | None = None, line: RGBColor | None = None, line_w: int = 0):
    shp = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shp.line.fill.background()
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w) if line_w else Pt(0.75)
    else:
        shp.line.fill.background()
    return shp

def txbox(slide, x, y, w, h, text: str, size: float = 14, bold=False, color: RGBColor = TEXT,
          align=PP_ALIGN.LEFT, font="Calibri", wrap=True, italic=False) -> object:
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    # Remove default text box margin for tighter control
    tb.left, tb.top, tb.width, tb.height = x, y, w, h
    return tb

def add_run(para, text: str, size: float, bold=False, color: RGBColor = TEXT,
            font="Calibri", italic=False):
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return run

def multiline_txbox(slide, x, y, w, h, lines: list[tuple], wrap=True):
    """lines = list of (text, size, bold, color, align, font, italic)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    first = True
    for item in lines:
        text, size, bold, color, align, font, italic = (item + (False,))[:7]
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
    return tb

def bullet_box(slide, x, y, w, h, items: list[str], size=12, color: RGBColor = TEXT,
               bullet_color: RGBColor = BRAND, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # Bullet using en-dash prefix for editability
        run = p.add_run()
        run.text = f"→  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = font
        p.space_before = Pt(3)
    return tb

# ── Background fills ──────────────────────────────────────────────────────────
def dark_bg(slide):
    """Full-bleed dark background."""
    rect(slide, 0, 0, W, H, fill=INK_950)

def light_bg(slide):
    """Warm off-white background."""
    rect(slide, 0, 0, W, H, fill=PAPER)

def top_bar(slide, color: RGBColor = BRAND, height=Inches(0.12)):
    """Brand-colored top accent bar."""
    rect(slide, 0, 0, W, height, fill=color)

def bottom_bar(slide, color: RGBColor = INK_900, height=Inches(0.18)):
    rect(slide, 0, H - height, W, height, fill=color)

def page_header(slide, section: str, dark=False):
    """Running header for content slides."""
    top_bar(slide, BRAND, Inches(0.06))
    fg = WHITE if dark else INK_500
    txbox(slide, Inches(0.4), Inches(0.07), Inches(3), Inches(0.2),
          "DOINg.AI — USER GUIDE", size=7, color=fg, font="Calibri")
    txbox(slide, 0, Inches(0.07), W - Inches(0.4), Inches(0.2),
          section.upper(), size=7, color=BRAND if not dark else WHITE,
          align=PP_ALIGN.RIGHT, font="Calibri")

def slide_number(slide, num: int, dark=False):
    """Page number chip bottom-left."""
    fg = BRAND
    rect(slide, Inches(0.3), H - Inches(0.36), Inches(0.32), Inches(0.22), fill=BRAND)
    txbox(slide, Inches(0.3), H - Inches(0.38), Inches(0.32), Inches(0.25),
          f"{num:02d}", size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(slide, Inches(0.7), H - Inches(0.36), Inches(4), Inches(0.2),
          "MAKE EVERY PROJECT VISIBLE  ·  DOING.AI", size=7, color=MUTED)

def section_chip(slide, x, y, text: str, bg=BRAND, fg=WHITE):
    """Small colored label chip."""
    w = Inches(1.8)
    h = Inches(0.22)
    rect(slide, x, y, w, h, fill=bg)
    txbox(slide, x + Inches(0.1), y, w - Inches(0.1), h,
          text.upper(), size=8, bold=True, color=fg, align=PP_ALIGN.LEFT)

def tip_box(slide, x, y, w, h, title: str, body: str):
    """Orange-bordered pro-tip callout box."""
    rect(slide, x, y, w, h, fill=RGBColor(0xFF,0xF4,0xEE), line=BRAND)
    rect(slide, x, y, Inches(0.07), h, fill=BRAND)   # left accent bar
    txbox(slide, x + Inches(0.16), y + Inches(0.1), w - Inches(0.2), Inches(0.22),
          f"💡  {title.upper()}", size=9, bold=True, color=BRAND)
    txbox(slide, x + Inches(0.16), y + Inches(0.32), w - Inches(0.2), h - Inches(0.36),
          body, size=10, color=TEXT, wrap=True)

def kpi_card(slide, x, y, w, h, label: str, value: str, sub: str = "", color: RGBColor = BRAND):
    rect(slide, x, y, w, h, fill=PAPER, line=INK_300, line_w=0.5)
    rect(slide, x, y, w, Inches(0.06), fill=color)
    txbox(slide, x + Inches(0.1), y + Inches(0.1), w - Inches(0.15), Inches(0.2),
          label.upper(), size=8, bold=True, color=MUTED)
    txbox(slide, x + Inches(0.1), y + Inches(0.3), w - Inches(0.15), Inches(0.45),
          value, size=28, bold=True, color=color)
    if sub:
        txbox(slide, x + Inches(0.1), y + h - Inches(0.25), w - Inches(0.15), Inches(0.22),
              sub, size=9, color=MUTED)

# ── Decorative grid pattern ───────────────────────────────────────────────────
def deco_grid(slide, x, y, cols=8, rows=6, cell=Inches(0.18), gap=Inches(0.04)):
    for r in range(rows):
        for c in range(cols):
            if (r + c) % 5 == 0:
                fill = BRAND
            elif (r + c) % 3 == 0:
                fill = INK_700
            else:
                continue
            rect(slide, x + c*(cell+gap), y - r*(cell+gap), cell, cell, fill=fill)

# ── SLIDE BUILDERS ────────────────────────────────────────────────────────────

def slide_cover(prs: Presentation, slide_n: int):
    s = blank_slide(prs)
    dark_bg(s)
    # Top brand bar
    rect(s, 0, 0, W, Inches(0.18), fill=BRAND)
    # Kicker
    txbox(s, Inches(0.55), Inches(0.25), Inches(9), Inches(0.22),
          f"EDITION {date.today().year}  ·  USER GUIDE  ·  ENGLISH",
          size=9, bold=True, color=BRAND)
    # Wordmark
    multiline_txbox(s, Inches(0.55), Inches(0.6), Inches(6.5), Inches(1.4), [
        ("DOINg", 72, True, WHITE, PP_ALIGN.LEFT, "Calibri"),
    ])
    # .AI suffix in brand color (overlaid — same position, wide box)
    txbox(s, Inches(3.35), Inches(0.6), Inches(3), Inches(1.4),
          ".AI", size=72, bold=True, color=BRAND)
    # Tagline
    txbox(s, Inches(0.55), Inches(1.85), Inches(7), Inches(0.35),
          "Make every project visible.", size=16, italic=True, color=INK_300)
    # Sub-headline
    txbox(s, Inches(0.55), Inches(2.25), Inches(6.2), Inches(0.8),
          "Your local-first AI Project Operations workspace.",
          size=22, bold=True, color=WHITE, wrap=True)

    # Three pillars
    pillars = [
        ("LOCAL-FIRST",    "Data on your machine.\nAlways available, never leaked."),
        ("AI-AUGMENTED",   "Local LLM enrichment.\nNo cloud calls."),
        ("TEAM-AWARE",     "SSE real-time sync.\nFour roles, one matrix."),
    ]
    pw = Inches(2.8)
    for i, (title, body) in enumerate(pillars):
        px = Inches(0.55) + i * (pw + Inches(0.15))
        rect(s, px, Inches(3.3), Inches(0.45), Inches(0.06), fill=BRAND)
        txbox(s, px, Inches(3.42), pw, Inches(0.24), title, size=10, bold=True, color=WHITE)
        txbox(s, px, Inches(3.68), pw, Inches(0.55), body, size=9, color=INK_300, wrap=True)

    # Decorative grid (right side)
    deco_grid(s, Inches(8.5), Inches(3.2), cols=7, rows=10, cell=Inches(0.13), gap=Inches(0.03))

    # Bottom strip
    rect(s, 0, H - Inches(0.36), W, Inches(0.36), fill=INK_900)
    txbox(s, Inches(0.4), H - Inches(0.3), Inches(4), Inches(0.24),
          f"V1.0  ·  {date.today().strftime('%B %Y').upper()}  ·  PUBLIC EDITION",
          size=8, color=INK_500)
    txbox(s, 0, H - Inches(0.3), W - Inches(0.4), Inches(0.24),
          "DOING.AI  ·  USER EDITION", size=8, color=WHITE, align=PP_ALIGN.RIGHT)
    return s

def slide_toc(prs, slide_n):
    s = blank_slide(prs)
    light_bg(s)
    page_header(s, "Contents")
    slide_number(s, slide_n)

    txbox(s, Inches(0.4), Inches(0.32), Inches(4), Inches(0.24),
          "TABLE OF CONTENTS", size=9, bold=True, color=BRAND)
    txbox(s, Inches(0.4), Inches(0.58), Inches(8), Inches(0.5),
          "What you will find inside", size=22, bold=True, color=TEXT)
    rect(s, Inches(0.4), Inches(1.08), Inches(0.55), Inches(0.05), fill=BRAND)

    entries = [
        ("01", "Welcome",             "Mission, principles, workspace tour."),
        ("02", "Project Management",  "Projects · Timeline · Risk · Communications."),
        ("03", "Catalogs",            "Technologies · Repos · Hackathons · MCP Hub · AI Agents."),
        ("04", "Collaboration",       "Working Groups and Weekly Check-in."),
        ("05", "Personal Tools",      "Smart ToDo and the Wish List."),
        ("06", "AI Companion",        "Assistant · Executive Insight · Command Palette."),
        ("07", "Notifications & You", "Alerts, profile, theme, self-service password change."),
        ("IT", "Data Feeds (IT)",     "Pipeline registry — admin + IT flag required."),
        ("A",  "Annex — Architecture","Local-first storage, real-time sync, LLM integration."),
        ("B",  "Annex — Permissions", "The full role × group matrix."),
        ("C",  "Annex — Shortcuts",   "Keyboard shortcuts and glossary."),
    ]
    col_break = 5  # First column: 5 items
    col_w = Inches(4.4)
    for i, (num, title, desc) in enumerate(entries):
        col = 0 if i < col_break else 1
        row = i if i < col_break else i - col_break
        x = Inches(0.4) + col * (col_w + Inches(0.5))
        y = Inches(1.22) + row * Inches(0.78)
        rect(s, x, y, Inches(0.32), Inches(0.32), fill=BRAND)
        txbox(s, x, y + Inches(0.04), Inches(0.32), Inches(0.24), num,
              size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txbox(s, x + Inches(0.4), y, Inches(3.7), Inches(0.22),
              title, size=11, bold=True, color=TEXT)
        txbox(s, x + Inches(0.4), y + Inches(0.22), Inches(3.7), Inches(0.2),
              desc, size=9, color=MUTED)
        rect(s, x + Inches(0.4), y + Inches(0.42), col_w - Inches(0.4), Inches(0.01), fill=INK_300)
    return s

def slide_divider(prs, num: str, title: str, kicker: str, snippet: str, color=BRAND):
    s = blank_slide(prs)
    dark_bg(s)
    # Giant numeral
    txbox(s, Inches(0.3), Inches(0.0), Inches(5), Inches(3.8),
          num, size=200, bold=True, color=BRAND, align=PP_ALIGN.LEFT)
    # Vertical accent line
    rect(s, Inches(5.8), Inches(0.4), Inches(0.04), Inches(4.8), fill=BRAND)
    # Kicker
    txbox(s, Inches(6.1), Inches(0.6), Inches(3.7), Inches(0.25),
          kicker.upper(), size=9, bold=True, color=BRAND)
    # Title
    txbox(s, Inches(6.1), Inches(0.9), Inches(3.7), Inches(0.65),
          title, size=28, bold=True, color=WHITE, wrap=True)
    # Snippet
    txbox(s, Inches(6.1), Inches(1.65), Inches(3.7), Inches(1.5),
          snippet, size=11, color=INK_300, wrap=True)
    # Deco grid bottom-left
    deco_grid(s, Inches(0.3), Inches(5.4), cols=12, rows=5, cell=Inches(0.13), gap=Inches(0.03))
    return s

def slide_welcome(prs, slide_n):
    s = blank_slide(prs)
    light_bg(s)
    page_header(s, "Welcome")
    slide_number(s, slide_n)
    section_chip(s, Inches(0.4), Inches(0.32), "01 · WELCOME")
    txbox(s, Inches(0.4), Inches(0.6), Inches(9), Inches(0.5),
          "What is DOINg.AI?", size=26, bold=True, color=TEXT)
    rect(s, Inches(0.4), Inches(1.1), Inches(0.55), Inches(0.05), fill=BRAND)

    txbox(s, Inches(0.4), Inches(1.22), Inches(9.2), Inches(1.2),
          "DOINg.AI is your team's local-first AI Project Operations workspace. It brings together "
          "every moving part of your innovation portfolio — projects, technologies, repositories, "
          "hackathons, MCP servers, AI agents, working groups, and personal todos — under one "
          "dark, keyboard-friendly canvas. Your data lives on your machine; nothing is shipped "
          "to the cloud unless you explicitly export it.",
          size=11, color=TEXT, wrap=True)

    # Three pillars
    pillars = [
        (BRAND, "LOCAL-FIRST",
         "localStorage + db.json on your network. Real-time SSE sync. Zero cloud dependency."),
        (BLUE,  "AI-AUGMENTED",
         "Plug Ollama, LM Studio, or n8n. Every prompt template editable. LLM stays on your network."),
        (GREEN, "TEAM-AWARE",
         "Four roles (admin / manager / contributor / viewer) shape every screen via a centralised permission matrix."),
    ]
    pw = Inches(2.9)
    py = Inches(2.6)
    for i, (col, title, body) in enumerate(pillars):
        px = Inches(0.4) + i * (pw + Inches(0.15))
        rect(s, px, py, pw, Inches(0.07), fill=col)
        txbox(s, px, py + Inches(0.12), pw, Inches(0.25), title, size=10, bold=True, color=TEXT)
        txbox(s, px, py + Inches(0.38), pw, Inches(0.6), body, size=9.5, color=MUTED, wrap=True)

    tip_box(s, Inches(0.4), Inches(3.9), Inches(9.2), Inches(1.45),
            "How to read this guide",
            "Sections 01–07 walk through each functional area. Annex A explains how the app is built; "
            "Annex B is the full permission matrix; Annex C is for keyboard-loving users. "
            "Every text element in this deck is editable — click any box and start typing.")
    return s

def slide_workspace_tour(prs, slide_n):
    s = blank_slide(prs)
    light_bg(s)
    page_header(s, "Workspace Tour")
    slide_number(s, slide_n)
    section_chip(s, Inches(0.4), Inches(0.32), "02 · WORKSPACE TOUR")
    txbox(s, Inches(0.4), Inches(0.6), Inches(9), Inches(0.45),
          "Three regions, one keyboard", size=24, bold=True, color=TEXT)
    rect(s, Inches(0.4), Inches(1.05), Inches(0.55), Inches(0.05), fill=BRAND)

    # Sidebar mock
    sx, sy, sw, sh = Inches(0.4), Inches(1.2), Inches(1.85), Inches(4.15)
    rect(s, sx, sy, sw, sh, fill=INK_900)
    txbox(s, sx + Inches(0.1), sy + Inches(0.08), Inches(0.7), Inches(0.22),
          "DOINg", size=10, bold=True, color=WHITE)
    txbox(s, sx + Inches(0.72), sy + Inches(0.08), Inches(0.55), Inches(0.22),
          ".AI", size=10, bold=True, color=BRAND)
    menu_items = ["Dashboard","Projects ●","Timeline","Risk Heatmap",
                  "Communications","Technologies","Repositories","Hackathons",
                  "MCP Hub","AI Agents","Working Groups","Weekly Check-in",
                  "Smart ToDo","Wish List"]
    for i, item in enumerate(menu_items):
        active = item.endswith("●")
        label = item.replace(" ●","")
        bg = BRAND if active else None
        if bg: rect(s, sx, sy + Inches(0.38) + i*Inches(0.26), sw, Inches(0.24), fill=bg)
        txbox(s, sx + Inches(0.1), sy + Inches(0.4) + i*Inches(0.26), sw - Inches(0.12),
              Inches(0.22), label.upper(), size=7.5, bold=True,
              color=WHITE if active else INK_300)
    # User card
    rect(s, sx, sy + sh - Inches(0.42), sw, Inches(0.42), fill=INK_700)
    rect(s, sx + Inches(0.1), sy + sh - Inches(0.36), Inches(0.25), Inches(0.25), fill=BRAND)
    txbox(s, sx + Inches(0.1), sy + sh - Inches(0.37), Inches(0.25), Inches(0.25),
          "JD", size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, sx + Inches(0.42), sy + sh - Inches(0.32), Inches(1.3), Inches(0.16),
          "JANE DOE", size=7.5, bold=True, color=WHITE)
    txbox(s, sx + Inches(0.42), sy + sh - Inches(0.17), Inches(1.3), Inches(0.15),
          "CONTRIBUTOR", size=6.5, color=BRAND)

    # Callouts
    callouts = [
        ("1", "Live sync indicator",
         "Green dot = backend reachable and pushing real-time SSE updates."),
        ("2", "Functional groups",
         "Sidebar items grouped by what they DO: management, catalogs, collaboration, tools."),
        ("3", "Active highlight",
         "Current screen highlighted in brand orange."),
        ("4", "Identity card",
         "Key icon → change password. Door icon → sign out."),
    ]
    for i, (num, title, body) in enumerate(callouts):
        cy = Inches(1.3) + i * Inches(1.02)
        rect(s, Inches(2.4), cy, Inches(0.3), Inches(0.3), fill=BRAND)
        txbox(s, Inches(2.4), cy + Inches(0.05), Inches(0.3), Inches(0.24),
              num, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txbox(s, Inches(2.8), cy, Inches(6.8), Inches(0.26),
              title, size=11, bold=True, color=TEXT)
        txbox(s, Inches(2.8), cy + Inches(0.26), Inches(6.8), Inches(0.4),
              body, size=9.5, color=MUTED, wrap=True)
        rect(s, Inches(2.4), cy + Inches(0.65), Inches(7.2), Inches(0.008), fill=INK_300)
    return s

def slide_feature(prs, slide_n, section_label, chip_color, chip_text,
                  title, tagline, what_text, bullet_items, tip_title, tip_body,
                  kpi_row=None):
    s = blank_slide(prs)
    light_bg(s)
    page_header(s, section_label)
    slide_number(s, slide_n)
    section_chip(s, Inches(0.4), Inches(0.32), chip_text, bg=chip_color)
    txbox(s, Inches(0.4), Inches(0.6), Inches(9), Inches(0.42),
          title, size=24, bold=True, color=TEXT)
    txbox(s, Inches(0.4), Inches(1.0), Inches(9), Inches(0.24),
          tagline, size=11, italic=True, color=MUTED)
    rect(s, Inches(0.4), Inches(1.22), Inches(0.55), Inches(0.04), fill=BRAND)

    # KPI row (optional)
    y_offset = Inches(0)
    if kpi_row:
        for i, (label, val, sub) in enumerate(kpi_row):
            kpi_card(s, Inches(0.4) + i * Inches(2.32), Inches(1.32),
                     Inches(2.15), Inches(0.9), label, val, sub, chip_color)
        y_offset = Inches(1.05)

    # What it does
    txbox(s, Inches(0.4), Inches(1.32) + y_offset, Inches(2), Inches(0.2),
          "WHAT IT DOES", size=8, bold=True, color=MUTED)
    txbox(s, Inches(0.4), Inches(1.52) + y_offset, Inches(5.6), Inches(1.1),
          what_text, size=10, color=TEXT, wrap=True)

    # Decorative feature block (right side)
    rect(s, Inches(6.2), Inches(1.25) + y_offset, Inches(3.4), Inches(1.5), fill=INK_950)
    rect(s, Inches(6.2), Inches(1.25) + y_offset, Inches(3.4), Inches(0.06), fill=chip_color)
    txbox(s, Inches(6.35), Inches(1.36) + y_offset, Inches(3.1), Inches(1.3),
          title.upper(), size=13, bold=True, color=chip_color, wrap=True)

    # Key actions
    txbox(s, Inches(0.4), Inches(2.7) + y_offset, Inches(2), Inches(0.2),
          "KEY ACTIONS", size=8, bold=True, color=MUTED)
    bullet_box(s, Inches(0.4), Inches(2.9) + y_offset, Inches(9.2),
               Inches(3.4) - y_offset - Inches(1.6), bullet_items, size=10.5)

    tip_box(s, Inches(0.4), Inches(4.6), Inches(9.2), Inches(0.88),
            tip_title, tip_body)
    return s

def slide_permission_matrix(prs, slide_n):
    s = blank_slide(prs)
    light_bg(s)
    page_header(s, "Annex B — Permissions")
    slide_number(s, slide_n)
    section_chip(s, Inches(0.4), Inches(0.32), "ANNEX B · ROLES & PERMISSIONS", bg=TEXT)
    txbox(s, Inches(0.4), Inches(0.6), Inches(9), Inches(0.42),
          "The Role × Group matrix", size=24, bold=True, color=TEXT)
    rect(s, Inches(0.4), Inches(1.02), Inches(0.55), Inches(0.04), fill=BRAND)

    # Matrix table
    headers = ["ROLE", "G1 — MGMT", "G2 — CATALOGS", "G4 — COLLAB", "PUBLIC"]
    col_ws = [1.2, 2.1, 2.3, 2.0, 1.5]
    rows = [
        ("admin",       "all view · edit · AI", "all view · edit · AI", "all",               "all"),
        ("manager",     "all view · edit · AI", "all view · edit · AI", "own view · edit · AI","all"),
        ("contributor", "own view · edit · AI", "all view, own edit",   "own view · edit · AI","all"),
        ("viewer",      "own view  (no edit)",  "all view, own edit",   "own view · edit",    "all"),
    ]
    x0, y0, rh = Inches(0.4), Inches(1.15), Inches(0.45)
    # Header row
    cx = x0
    rect(s, x0, y0, sum(Inches(w) for w in col_ws), rh, fill=INK_950)
    for i, hdr in enumerate(headers):
        txbox(s, cx + Inches(0.08), y0 + Inches(0.12), Inches(col_ws[i]) - Inches(0.1), Inches(0.22),
              hdr, size=9, bold=True, color=WHITE)
        cx += Inches(col_ws[i])
    # Data rows
    zebra = [PAPER, RGBColor(0xF0,0xEF,0xE9)]
    for ri, row in enumerate(rows):
        ry = y0 + (ri+1)*rh
        cx = x0
        rect(s, x0, ry, sum(Inches(w) for w in col_ws), rh, fill=zebra[ri%2])
        for ci, cell in enumerate(row):
            c = TEXT if ci == 0 else MUTED
            txbox(s, cx + Inches(0.08), ry + Inches(0.13), Inches(col_ws[ci]) - Inches(0.1),
                  Inches(0.3), cell, size=9, bold=(ci==0), color=c)
            cx += Inches(col_ws[ci])
        rect(s, x0, ry + rh - Inches(0.01), sum(Inches(w) for w in col_ws), Inches(0.01), fill=INK_300)

    # Group glossary
    y_g = y0 + 5 * rh + Inches(0.2)
    txbox(s, Inches(0.4), y_g, Inches(2), Inches(0.22), "GROUP GLOSSARY", size=8, bold=True, color=MUTED)
    glossary = [
        ("G1", "Projects, Timeline, Risk, Communications."),
        ("G2", "Technologies, Repos, Hackathons, MCP Hub, AI Agents."),
        ("G4", "Working Groups, Weekly Check-in."),
        ("Public", "Dashboard, Smart ToDo, Wish List, User Guide."),
    ]
    y_g += Inches(0.28)
    col_w_g = Inches(4.4)
    for i, (tag, desc) in enumerate(glossary):
        col = 0 if i < 2 else 1
        row_g = i if i < 2 else i - 2
        gx = Inches(0.4) + col * col_w_g
        gy = y_g + row_g * Inches(0.35)
        rect(s, gx, gy, Inches(0.42), Inches(0.25), fill=BRAND)
        txbox(s, gx, gy + Inches(0.04), Inches(0.42), Inches(0.2),
              tag, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txbox(s, gx + Inches(0.5), gy, Inches(col_w_g - Inches(0.55)), Inches(0.28),
              desc, size=9, color=TEXT)
    return s

def slide_architecture(prs, slide_n):
    s = blank_slide(prs)
    light_bg(s)
    page_header(s, "Annex A — Architecture")
    slide_number(s, slide_n)
    section_chip(s, Inches(0.4), Inches(0.32), "ANNEX A · TECHNICAL ARCHITECTURE", bg=TEXT)
    txbox(s, Inches(0.4), Inches(0.6), Inches(9), Inches(0.42),
          "How the workspace fits together", size=24, bold=True, color=TEXT)
    rect(s, Inches(0.4), Inches(1.02), Inches(0.55), Inches(0.04), fill=BRAND)
    txbox(s, Inches(0.4), Inches(1.1), Inches(9.2), Inches(0.32),
          "Local-first, three-tier. Everything below your browser runs on hardware YOU control — no third-party SaaS in the critical path.",
          size=10, color=MUTED, wrap=True)

    # Architecture boxes
    tiers = [
        (BRAND,   "BROWSER\nReact 19 + TypeScript",
         ["UI Components","Permission filter","localStorage cache","SSE event listener"]),
        (INK_700, "BACKEND\nFastAPI (Python)",
         ["/api/data (CRUD)","SSE broadcast (/api/events)","Proxy endpoints","Version check"]),
        (BRAND,   "STORAGE + AI\nLocal stack",
         ["db.json on disk","localStorage mirror","Local LLM (Ollama…)","External proxies"]),
    ]
    bw = Inches(2.8)
    bh = Inches(2.5)
    by = Inches(1.55)
    for i, (col, title, items) in enumerate(tiers):
        bx = Inches(0.4) + i * (bw + Inches(0.8))
        rect(s, bx, by, bw, bh, fill=INK_950)
        rect(s, bx, by, bw, Inches(0.06), fill=col)
        txbox(s, bx + Inches(0.12), by + Inches(0.1), bw - Inches(0.2), Inches(0.5),
              title, size=9, bold=True, color=col, wrap=True)
        for j, item in enumerate(items):
            iy = by + Inches(0.65) + j * Inches(0.43)
            item_col = col if j == 0 else INK_700
            rect(s, bx + Inches(0.12), iy, bw - Inches(0.24), Inches(0.33), fill=item_col)
            txbox(s, bx + Inches(0.2), iy + Inches(0.08), bw - Inches(0.32), Inches(0.22),
                  item, size=8.5, bold=True, color=WHITE)
        # Arrow to next tier
        if i < 2:
            ax = bx + bw + Inches(0.15)
            rect(s, ax, by + bh//2 - Inches(0.02), Inches(0.55), Inches(0.04), fill=BRAND)
            txbox(s, ax + Inches(0.28), by + bh//2 - Inches(0.2), Inches(0.28), Inches(0.4),
                  "▶", size=10, bold=True, color=BRAND)

    # Two narrative blocks below
    txbox(s, Inches(0.4), Inches(4.22), Inches(2), Inches(0.2),
          "WRITE PATH", size=8, bold=True, color=MUTED)
    txbox(s, Inches(0.4), Inches(4.42), Inches(4.5), Inches(0.9),
          "You edit a field → React updates state → writes to localStorage IMMEDIATELY → "
          "debounces a POST to /api/data (400ms). X-Base-Version for optimistic concurrency.",
          size=9, color=TEXT, wrap=True)
    txbox(s, Inches(5.1), Inches(4.22), Inches(2), Inches(0.2),
          "READ / SYNC PATH", size=8, bold=True, color=MUTED)
    txbox(s, Inches(5.1), Inches(4.42), Inches(4.5), Inches(0.9),
          "Backend pushes 'data_updated' over SSE the instant any client writes. Other tabs "
          "refresh in < 100ms. 30s polling acts as a safety net. Pending edits always win.",
          size=9, color=TEXT, wrap=True)
    return s

def slide_shortcuts(prs, slide_n):
    s = blank_slide(prs)
    light_bg(s)
    page_header(s, "Annex C — Shortcuts")
    slide_number(s, slide_n)
    section_chip(s, Inches(0.4), Inches(0.32), "ANNEX C · KEYBOARD + GLOSSARY", bg=TEXT)
    txbox(s, Inches(0.4), Inches(0.6), Inches(9), Inches(0.4),
          "Keyboard shortcuts & glossary", size=24, bold=True, color=TEXT)
    rect(s, Inches(0.4), Inches(1.0), Inches(0.55), Inches(0.04), fill=BRAND)

    shortcuts = [
        ("⌘ K  /  Ctrl K",     "Open the Command Palette (fuzzy navigation + project jump)."),
        ("Esc",                 "Close any modal or dialog."),
        ("Enter",               "Confirm in dialogs; send in the assistant chat."),
        ("Shift + Enter",       "New line in the assistant chat (without sending)."),
        ("Tab  /  Shift+Tab",   "Move between form fields. Fully keyboard-navigable."),
    ]
    txbox(s, Inches(0.4), Inches(1.1), Inches(3), Inches(0.22),
          "KEYBOARD SHORTCUTS", size=8, bold=True, color=MUTED)
    for i, (keys, desc) in enumerate(shortcuts):
        ky = Inches(1.35) + i * Inches(0.5)
        rect(s, Inches(0.4), ky, Inches(2.2), Inches(0.32), fill=INK_950)
        txbox(s, Inches(0.5), ky + Inches(0.08), Inches(2.1), Inches(0.22),
              keys, size=10, bold=True, color=BRAND, font="Consolas")
        txbox(s, Inches(2.72), ky + Inches(0.08), Inches(4.5), Inches(0.22),
              desc, size=9.5, color=TEXT)

    txbox(s, Inches(0.4), Inches(3.9), Inches(2), Inches(0.22),
          "GLOSSARY", size=8, bold=True, color=MUTED)
    glossary = [
        ("MCP",           "Model Context Protocol — the standard for LLMs calling external tools."),
        ("Agent",         "Autonomous AI program with a goal, tools, and prompts."),
        ("Family",        "Label that bundles related projects or MCPs together."),
        ("SSE",           "Server-Sent Events — one-way push channel from backend to browser."),
        ("Silver copy",   "Clean, query-ready copy of source data in a data warehouse."),
        ("Working Group", "Standing cross-functional squad with tasks, notes, and decisions."),
        ("Wish",          "Documented need (project / MCP / agent / feature) routed to a sponsor."),
    ]
    for i, (term, definition) in enumerate(glossary):
        col = 0 if i < 4 else 1
        row = i if i < 4 else i - 4
        gx = Inches(0.4) + col * Inches(4.8)
        gy = Inches(4.15) + row * Inches(0.35)
        txbox(s, gx, gy, Inches(1.3), Inches(0.28), term, size=9.5, bold=True, color=BRAND)
        txbox(s, gx + Inches(1.35), gy, Inches(3.3), Inches(0.28), definition, size=9, color=TEXT, wrap=True)
    return s

def slide_ai_companion(prs, slide_n):
    s = blank_slide(prs)
    light_bg(s)
    page_header(s, "AI Companion")
    slide_number(s, slide_n)
    section_chip(s, Inches(0.4), Inches(0.32), "07 · AI COMPANION  ·  ⌘K")
    txbox(s, Inches(0.4), Inches(0.6), Inches(9), Inches(0.42),
          "Your AI companion — locally hosted", size=24, bold=True, color=TEXT)
    txbox(s, Inches(0.4), Inches(1.0), Inches(9), Inches(0.24),
          "All AI surfaces route to YOUR configured local LLM. Nothing leaves your network.", size=11, italic=True, color=MUTED)
    rect(s, Inches(0.4), Inches(1.22), Inches(0.55), Inches(0.04), fill=BRAND)

    items = [
        (BRAND, "DOINg Assistant",
         "Open from the top bar. Ask 'what slipped this week?', 'summarise the AI Efficiency family', "
         "'draft a stakeholder email for project X'. The assistant has your filtered portfolio in context.",
         ["Full portfolio context in the chat","Markdown responses with references","Works with any OpenAI-compatible endpoint"]),
        (BLUE, "AI Executive Insight",
         "Click the orange 'AI EXECUTIVE INSIGHT' button in the sidebar. Get a structured one-paragraph "
         "portfolio brief fit for an Exco deck — one click, no configuration.",
         ["Portfolio KPIs automatically included","English only (by default — prompt is editable)","Copy with one click"]),
        (GREEN, "Command Palette  ⌘K",
         "Press ⌘K. Fuzzy-jump to any view, search a project by name, or trigger a recently-used action. "
         "Faster than your cursor will ever be.",
         ["Fuzzy search across all names and tags","Jump directly to a project card","⌘K from anywhere in the app"]),
    ]
    cw = Inches(3.0)
    for i, (col, title, body, bullets) in enumerate(items):
        cx = Inches(0.4) + i * (cw + Inches(0.14))
        # Color top bar
        rect(s, cx, Inches(1.35), cw, Inches(0.06), fill=col)
        rect(s, cx, Inches(1.41), cw, Inches(1.9), fill=RGBColor(0xF9,0xFB,0xFF) if col != BRAND else RGBColor(0xFF,0xF8,0xF5))
        txbox(s, cx + Inches(0.12), Inches(1.46), cw - Inches(0.18), Inches(0.28),
              title, size=12, bold=True, color=col)
        txbox(s, cx + Inches(0.12), Inches(1.76), cw - Inches(0.18), Inches(1.3),
              body, size=9, color=TEXT, wrap=True)
        for j, b in enumerate(bullets):
            by = Inches(3.4) + j * Inches(0.34)
            rect(s, cx + Inches(0.12), by + Inches(0.07), Inches(0.07), Inches(0.07), fill=col)
            txbox(s, cx + Inches(0.26), by, cw - Inches(0.3), Inches(0.3),
                  b, size=9, color=MUTED)

    tip_box(s, Inches(0.4), Inches(4.6), Inches(9.2), Inches(0.88),
            "When the assistant gets stuck",
            "Open Settings → Local LLM (ask an admin) to verify the connection. The assistant shows a "
            "single error toast and stays usable for browsing the rest of the app while you debug.")
    return s

def slide_notifications_profile(prs, slide_n):
    s = blank_slide(prs)
    light_bg(s)
    page_header(s, "Notifications & Profile")
    slide_number(s, slide_n)
    section_chip(s, Inches(0.4), Inches(0.32), "07 · NOTIFICATIONS & YOU")
    txbox(s, Inches(0.4), Inches(0.6), Inches(9), Inches(0.42),
          "Stay on track, stay you", size=24, bold=True, color=TEXT)
    rect(s, Inches(0.4), Inches(1.02), Inches(0.55), Inches(0.04), fill=BRAND)

    # Two columns
    cols = [
        ("NOTIFICATIONS", "🔔 The Bell", [
            "Project you own/are member of → silent > 10 days",
            "Task assigned to you → past its ETA",
            "Weekly check-in → overdue by more than 8 days",
            "Click the bell → dismiss alerts, mark all read",
        ]),
        ("PROFILE & SECURITY", "🔑 Change Password", [
            "Click the KEY icon next to your name in the sidebar",
            "Enter current password + new password (min 8 chars)",
            "Strength bar guides you to 'Strong'",
            "Passwords hashed with PBKDF2-SHA256 — never plain text",
        ]),
    ]
    col_w = Inches(4.5)
    for i, (label, title, bullets) in enumerate(cols):
        cx = Inches(0.4) + i * (col_w + Inches(0.6))
        txbox(s, cx, Inches(1.15), col_w, Inches(0.22), label, size=8, bold=True, color=MUTED)
        txbox(s, cx, Inches(1.4), col_w, Inches(0.32), title, size=14, bold=True, color=TEXT)
        bullet_box(s, cx, Inches(1.78), col_w, Inches(2.0), bullets, size=10, color=TEXT)

    tip_box(s, Inches(0.4), Inches(4.6), Inches(9.2), Inches(0.88),
            "Forgot your password?",
            "Ask an admin to use Settings → Security to send you a temporary password. "
            "You will be forced to choose a new one at the next sign-in. "
            "Admins can also toggle the 'IT' flag to give you access to the Data Feeds section.")
    return s

def slide_tips(prs, slide_n):
    s = blank_slide(prs)
    light_bg(s)
    page_header(s, "Tips & Best Practice")
    slide_number(s, slide_n)
    section_chip(s, Inches(0.4), Inches(0.32), "GROW INTO IT")
    txbox(s, Inches(0.4), Inches(0.6), Inches(9), Inches(0.42),
          "Habits that compound", size=24, bold=True, color=TEXT)
    rect(s, Inches(0.4), Inches(1.02), Inches(0.55), Inches(0.04), fill=BRAND)

    tips = [
        ("Update your check-in every Monday",
         "Two minutes. Feeds team rituals, the dashboard, and AI portfolio briefs. Skip → alert appears."),
        ("Live in the command palette  ⌘K",
         "Three letters. Jump. Faster than any menu — always."),
        ("Use AI as a brain-dump partner",
         "Paste raw meeting notes or a URL into a description; ask 'what should I prepare for Friday?' It has read access."),
        ("Tag aggressively, file lightly",
         "Two short tags beat one paragraph of taxonomy. Dashboard and search both lean on tags."),
        ("Use the Wish List instead of side-channel asks",
         "Need an MCP for system X? Drop it in the Wish List with a sponsor. Conversations live next to the request."),
        ("Treat the dashboard like a daily standup mirror",
         "Five minutes a morning. Anything red → act. Anything green → ignore."),
    ]
    col_w = Inches(4.4)
    for i, (title, body) in enumerate(tips):
        col = 0 if i < 3 else 1
        row = i if i < 3 else i - 3
        tx = Inches(0.4) + col * (col_w + Inches(0.6))
        ty = Inches(1.2) + row * Inches(1.15)
        rect(s, tx, ty, Inches(0.07), Inches(0.95), fill=BRAND)
        txbox(s, tx + Inches(0.17), ty, col_w - Inches(0.17), Inches(0.3),
              title, size=11, bold=True, color=TEXT)
        txbox(s, tx + Inches(0.17), ty + Inches(0.3), col_w - Inches(0.17), Inches(0.6),
              body, size=9.5, color=MUTED, wrap=True)
    return s

def slide_data_feeds(prs, slide_n, TEAL):
    s = blank_slide(prs)
    light_bg(s)
    page_header(s, "Data Platform")
    slide_number(s, slide_n)

    # Restricted-access banner at top
    rect(s, Inches(0.4), Inches(0.28), Inches(9.2), Inches(0.26), fill=TEAL)
    txbox(s, Inches(0.55), Inches(0.29), Inches(8.8), Inches(0.22),
          "🔒  RESTRICTED ACCESS — visible to Admin and users with the IT flag activated in Settings → Security",
          size=8, bold=True, color=WHITE)

    section_chip(s, Inches(0.4), Inches(0.62), "G5 · DATA PLATFORM", bg=TEAL)
    txbox(s, Inches(0.4), Inches(0.9), Inches(9), Inches(0.4),
          "Data Feeds", size=24, bold=True, color=TEXT)
    txbox(s, Inches(0.4), Inches(1.28), Inches(9), Inches(0.22),
          "The silver-copy & data-warehouse pipeline registry.", size=11, italic=True, color=MUTED)
    rect(s, Inches(0.4), Inches(1.48), Inches(0.55), Inches(0.04), fill=TEAL)

    # What it does
    txbox(s, Inches(0.4), Inches(1.58), Inches(1.8), Inches(0.2),
          "WHAT IT DOES", size=8, bold=True, color=MUTED)
    txbox(s, Inches(0.4), Inches(1.78), Inches(5.6), Inches(0.85),
          "Centralised registry of every data pipeline that feeds your silver copies and data warehouses. "
          "Each feed documents the source platform, data type, frequency, source → destination mapping, "
          "production date, JIRA reference, cost in Man-Days, and the PM and developer responsible.",
          size=10, color=TEXT, wrap=True)

    # Feature visual block (right side)
    rect(s, Inches(6.2), Inches(1.52), Inches(3.4), Inches(1.2), fill=INK_950)
    rect(s, Inches(6.2), Inches(1.52), Inches(3.4), Inches(0.06), fill=TEAL)
    txbox(s, Inches(6.35), Inches(1.63), Inches(3.1), Inches(1.0),
          "DATA FEEDS", size=13, bold=True, color=TEAL, wrap=True)

    # Field grid — 2 columns
    fields_left = [
        ("Platform name",    "Source system (ex: Salesforce, SAP, Oracle)"),
        ("Data type",        "Nature of the data (CRM events, Finance ledger…)"),
        ("Frequency",        "realtime / hourly / daily / weekly / monthly / on_demand"),
        ("Source",           "Origin endpoint — DB, API, file, streaming"),
        ("Destination",      "Silver copy table or DWH schema"),
        ("Status",           "planned → in_dev → uat → production → deprecated"),
    ]
    fields_right = [
        ("Prod date",        "Date the feed went live in production"),
        ("JIRA ref",         "Reference ticket for traceability"),
        ("Cost (Man-Days)",  "Estimated delivery effort"),
        ("ETA",              "Expected go-live if not yet in production"),
        ("Project Manager",  "PM name (free text)"),
        ("Developer",        "Dev name + presentation link"),
    ]
    col_w = Inches(4.45)
    for col_idx, fields in enumerate([fields_left, fields_right]):
        cx = Inches(0.4) + col_idx * (col_w + Inches(0.3))
        for row_idx, (label, desc) in enumerate(fields):
            fy = Inches(2.72) + row_idx * Inches(0.32)
            rect(s, cx, fy + Inches(0.06), Inches(0.08), Inches(0.08), fill=TEAL)
            txbox(s, cx + Inches(0.16), fy, Inches(1.35), Inches(0.28),
                  label, size=9, bold=True, color=TEXT)
            txbox(s, cx + Inches(1.52), fy, col_w - Inches(1.55), Inches(0.28),
                  desc, size=9, color=MUTED, wrap=True)

    tip_box(s, Inches(0.4), Inches(4.63), Inches(9.2), Inches(0.82),
            "Who can see this section?",
            "Data Feeds is visible to Admin users AND any user whose profile has the IT flag enabled. "
            "An admin can toggle the IT flag per user in Settings → Security. "
            "Use this section to maintain a full inventory of your data pipelines with clear ownership and SLA context.")
    return s


def slide_back_cover(prs):
    s = blank_slide(prs)
    dark_bg(s)
    rect(s, 0, H - Inches(0.22), W, Inches(0.22), fill=BRAND)
    txbox(s, Inches(0.55), Inches(0.6), Inches(7), Inches(1.1),
          "Local. Fast. AI-augmented.", size=40, bold=True, color=WHITE, wrap=True)
    txbox(s, Inches(0.55), Inches(1.72), Inches(4), Inches(0.72),
          "Yours.", size=40, bold=True, color=BRAND)
    rect(s, Inches(0.55), Inches(2.5), Inches(0.9), Inches(0.06), fill=BRAND)
    txbox(s, Inches(0.55), Inches(2.65), Inches(7.5), Inches(1.2),
          "You now know what DOINg.AI can do for your team. Keep this deck handy for new joiners, "
          "or jump back into the in-app User Guide whenever you need a quick reminder — "
          "the content there is always the live source of truth.",
          size=12, color=INK_300, wrap=True)
    txbox(s, Inches(0.55), Inches(4.6), Inches(8), Inches(0.3),
          f"EDITION {date.today().year}  ·  V1.0  ·  GENERATED {date.today().isoformat()}  ·  PUBLIC EDITION  ·  ADMIN FEATURES NOT INCLUDED",
          size=8, color=INK_500)
    deco_grid(s, Inches(8.2), Inches(4.8), cols=8, rows=20, cell=Inches(0.13), gap=Inches(0.03))
    return s

# ── MAIN ─────────────────────────────────────────────────────────────────────
def build(output_path: str):
    prs = new_prs()
    n = 1

    slide_cover(prs, n); n += 1
    slide_toc(prs, n); n += 1

    # 01 — Welcome
    slide_divider(prs, "01", "WELCOME", "YOUR LOCAL-FIRST WORKSPACE",
        "A quick orientation to what DOINg.AI is, the three pillars that shape every screen, and how to use this deck."); n += 1
    slide_welcome(prs, n); n += 1

    # 02 — Workspace Tour
    slide_divider(prs, "02", "WORKSPACE TOUR", "NAVIGATE WITHOUT FRICTION",
        "The three regions of the workspace, how the sidebar is organised, and the always-on AI surfaces."); n += 1
    slide_workspace_tour(prs, n); n += 1

    # 03 — Project Management
    slide_divider(prs, "03", "PROJECT MANAGEMENT", "WHERE THE WORK LANDS",
        "Projects, Timeline, Risk Heatmap, Communications — the four screens that drive day-to-day delivery."); n += 1

    slide_feature(prs, n, "Project Management", BRAND, "G1 · PROJECT MANAGEMENT",
        "Projects", "The heart of your portfolio.",
        "Rich record: description, context, manager, members, tasks, milestones, tags, technologies, "
        "repositories, audit log. Status, dev stage, and innovation maturity are first-class.",
        ["Create from scratch, from a template, or by importing raw notes via AI.",
         "Assign a manager and members — the assignment opens visibility for them.",
         "Add tasks with priority, ETA, weight, and dependencies — they roll up to the timeline.",
         "Flag projects as 'important' or 'archived'; filters everywhere honour those flags."],
        "Pro tip — use the Family field",
        "Group projects under a Family (e.g. 'AI Efficiency'). The Dashboard breaks down KPIs by family, "
        "and the Wish List can reference them directly."); n += 1

    slide_feature(prs, n, "Project Management", BRAND, "G1 · PROJECT MANAGEMENT",
        "Timeline", "A Gantt across the whole portfolio.",
        "Every project becomes a bar from start date to deadline. Milestones appear as diamonds; "
        "tasks are nested rows. Filter by family, status, manager, or tag to focus, or zoom out to spot "
        "deadline clashes a quarter ahead.",
        ["Drag a bar to shift a project window — the change persists immediately.",
         "Click a milestone diamond to mark it done or edit its date.",
         "Use the today line to spot what should land this week.",
         "Hover any bar for a 2-line tooltip with project status and manager."],
        "Pro tip — print this view",
        "The browser print dialog (⌘P) renders the timeline as a landscape PDF — perfect for a "
        "quarterly steering deck."); n += 1

    slide_feature(prs, n, "Project Management", BRAND, "G1 · PROJECT MANAGEMENT",
        "Risk Heatmap", "Probability × Impact, at a glance.",
        "Every project carries a risk profile. The Heatmap renders the portfolio on a 2D grid: "
        "probability on one axis, impact on the other. Hot zones glow brand orange so the highest-priority "
        "items pop immediately.",
        ["Adjust a project's risk score inline — heatmap updates instantly.",
         "Filter by family to compare risk profiles across initiatives.",
         "Use the AI mitigation generator to draft a starter mitigation plan."],
        "Pro tip — review weekly",
        "Treat the heatmap like a pre-meeting check. Anything in the top-right quadrant deserves "
        "a 60-second talk in your next stand-up — even if everything else is green."); n += 1

    slide_feature(prs, n, "Project Management", BRAND, "G1 · PROJECT MANAGEMENT",
        "Communications", "Drafts, newsletters, weeklies — in one place.",
        "Compose weekly status updates, monthly newsletters, Exco briefs, or info posts. Target "
        "mailing lists, attach project context, and let the LLM draft the body from your bullet points.",
        ["Pick a template (weekly, newsletter, exco, info) — each has its own prompt.",
         "Attach related projects so the AI has full context when drafting.",
         "Choose mailing lists; recipient count is computed live.",
         "Save as draft or mark as sent — the audit log preserves the trail."],
        "Pro tip — start from a brief",
        "On a project page, click 'Generate exec brief'. Copy that into a new Communication and the AI "
        "will tighten it into a newsletter — no blank-page anxiety."); n += 1

    # 04 — Catalogs
    slide_divider(prs, "04", "CATALOGS", "SHARED ASSETS · LIVING INVENTORY",
        "Technologies, repositories, hackathons, MCP servers, AI agents — the building blocks the whole team can reach for."); n += 1

    slide_feature(prs, n, "Catalogs", BLUE, "PUBLIC · TECHNOLOGIES",
        "Technologies", "A living radar of your stack.",
        "Catalogue every framework, library, database, language, or service. Each entry carries a layer "
        "(frontend / backend / data / ML-AI) and a maturity status (evaluating / adopted / deprecated / hold).",
        ["Add a tech with name, category, layer, and maturity in one click.",
         "Browse by layer or maturity; filter to 'evaluating' to see your innovation funnel.",
         "Link a project to a tech — it appears on both pages.",
         "Use AI enrichment to auto-fill description, version, license, and URL."],
        "Pro tip — quarterly tech review",
        "Once a quarter, scan the 'Evaluating' tier with the team. Promote what stuck, deprecate what "
        "you stopped using. The catalog should reflect reality, not aspirations."); n += 1

    slide_feature(prs, n, "Catalogs", BLUE, "G2 · CODE REPOSITORIES",
        "Repositories", "Source-control, mapped to projects.",
        "Directory of every repo (Bitbucket, GitHub, GitLab, Azure) your team owns or contributes to. "
        "Visibility, provider, language, and project links are all first-class.",
        ["Add a repo manually or bulk-import from your provider with a single token.",
         "Link a repo to one or more projects — appears on both records.",
         "Use the visibility chip to spot private repos at a glance.",
         "Search by language to find your Python or TypeScript estate in a click."],
        "Pro tip — keep it accurate",
        "When you spin up a new repo, add it here the same day. A 30-second discipline; the catalog "
        "stays trustworthy for everyone else."); n += 1

    slide_feature(prs, n, "Catalogs", BLUE, "G2 · HACKATHONS",
        "Hackathons", "From kick-off to retrospective in one place.",
        "Run hackathons end-to-end: theme, objective, challenge, dates, location, participants "
        "(devs, SMEs, experts, PMs), docs (briefs, guides, results), and a live message thread.",
        ["Plan with status (upcoming / active / completed / archived) and dates.",
         "Add participants with roles and teams; export the roster.",
         "Drop briefs and guides as docs; participants can author and edit.",
         "When done, generate the AI synthesis for a polished retrospective."],
        "Pro tip — link to a project",
        "If a hackathon produces something real, create a project AND a Wish List entry that references both. "
        "The lineage is captured forever in your audit log."); n += 1

    slide_feature(prs, n, "Catalogs", BLUE, "G2 · MCP HUB",
        "MCP Hub", "Your Model Context Protocol catalogue.",
        "Every MCP server (bridges that let LLMs call your internal systems) lives here. Each card "
        "carries description, tools, deploy status (dev / UAT / production), family, scope, and an "
        "AI-driven code analysis with security/perf recommendations.",
        ["Add a server declaratively or by URL — the Hub inventories its tools.",
         "Group servers in families with a custom colour for cross-portfolio visibility.",
         "Run AI code analysis to surface security/perf issues — track fixes inline.",
         "Export a PDF booklet of the whole Hub for an Exco committee."],
        "Pro tip — share with sponsors",
        "The PDF booklet is purpose-built for non-technical stakeholders. One click, polished output, "
        "brand-consistent. Drop it in Slack or attach to a steering deck."); n += 1

    slide_feature(prs, n, "Catalogs", BLUE, "G2 · AI AGENTS",
        "AI Agents", "Track internal agents from idea to production.",
        "A purpose-built tracker that adds agent-specific fields: model, framework (LangChain, LlamaIndex, "
        "CrewAI, AutoGen, n8n, custom), system prompt, knowledge bases, tools, and release/eval history.",
        ["Track agents through 6 stages: idea → design → dev → testing → production → deprecated.",
         "Link an agent to the MCP servers it depends on.",
         "Capture monthly invocations, cost, success rate for ops dashboards.",
         "Use families and tags to group by business domain."],
        "Pro tip — log every eval",
        "After each regression run, add an evaluation record. Three months in, you have a real time series — "
        "quality and cost trends you can show to leadership."); n += 1

    # IT Platform section — Data Feeds (G5: admin + IT flag)
    TEAL = RGBColor(0x02, 0x80, 0x90)
    slide_divider(prs, "IT", "DATA PLATFORM", "G5 · ADMIN + IT FLAG REQUIRED",
        "Data Feeds — the silver-copy pipeline registry. Visible to admin and users with the IT flag activated by an admin.",
        color=TEAL); n += 1
    slide_data_feeds(prs, n, TEAL); n += 1

    # 05 — Collaboration
    slide_divider(prs, "05", "COLLABORATION", "WORKING TOGETHER · ASYNC-FIRST",
        "Working Groups for standing squads and Weekly Check-in to keep the team rhythm."); n += 1

    slide_feature(prs, n, "Collaboration", GREEN, "G4 · WORKING GROUPS",
        "Working Groups", "Standing squads with memory.",
        "A long-running cross-functional squad with an objective, roster (lead / contributor / observer), "
        "task board (Kanban with priorities and checklists), and meeting notes with decisions and action items.",
        ["Create a group; invite members with role (lead / contributor / observer).",
         "Run sessions: capture agenda, notes, decisions, and action items.",
         "Action items carry forward — open ones surface in the next meeting prep.",
         "Track tasks in a focused board (todo / doing / review / done)."],
        "Pro tip — close the loop",
        "At every session, do a 2-minute 'carried action' review. The system shows you what was opened "
        "last time. Nothing falls through the cracks."); n += 1

    slide_feature(prs, n, "Collaboration", GREEN, "G4 · WEEKLY CHECK-IN",
        "Weekly Check-in", "Two minutes, every Monday.",
        "A short structured form: accomplishments, blockers, next steps, mood (green / amber / red). "
        "Submit it; it shows up on the team dashboard and feeds the AI weekly digest.",
        ["Reuse last week as a starting point — only edit what changed.",
         "Set mood with one click; the dashboard shows team weather at a glance.",
         "Read peers' check-ins in the dashboard digest for async stand-up vibes.",
         "Skipping a week triggers a friendly nudge in your notifications."],
        "Pro tip — block 2 minutes Monday morning",
        "Same time every week. After a month, it's automatic and your manager has constant, "
        "low-overhead visibility."); n += 1

    # 06 — Personal Tools
    slide_divider(prs, "06", "PERSONAL TOOLS", "YOUR SPACE INSIDE THE WORKSPACE",
        "Smart ToDo for everything that doesn't deserve a project, and the Wish List for what you want from the team."); n += 1

    slide_feature(prs, n, "Personal Tools", AMBER, "PUBLIC · SMART TODO",
        "Smart ToDo", "A personal task tracker with AI-grade structure.",
        "For anything that isn't a full project but still needs tracking. Each item carries requester, "
        "sponsor, priority, Eisenhower quadrant, energy required, estimated duration, and planning horizon.",
        ["Create todos in one line; the AI parses dates and priority.",
         "Plan by horizon (today / week / month / quarter / year) — focus mode hides the rest.",
         "Track time spent vs estimated — useful for personal calibration.",
         "Mark recurring tasks; new instances roll forward automatically."],
        "Pro tip — sponsor field matters",
        "When a todo has a sponsor, you can resurface 'what does this person need from me' in two clicks. "
        "Big psychological win when you start a tough conversation."); n += 1

    slide_feature(prs, n, "Personal Tools", AMBER, "PUBLIC · WISH LIST",
        "Wish List", "Document your needs. Find your sponsor.",
        "A collective backlog where anyone can submit a need — a new project, an MCP server, an AI agent, "
        "a data feed, a feature, an integration. Each wish carries sponsor, dates, justification, and a comment thread.",
        ["Submit a wish in 30 seconds — title + description + category is enough to start.",
         "Pick an internal sponsor from your colleagues, or free-text an external name.",
         "Cross-link to existing projects, MCPs, or agents — dropdown is category-aware.",
         "Comment on anyone's wish; track status (submitted → accepted → in progress → done)."],
        "Pro tip — sponsor before you submit",
        "A quick 5-minute chat with your potential sponsor before you submit dramatically increases acceptance. "
        "The Wish List is the contract; the conversation is the trust."); n += 1

    # 07 — AI Companion
    slide_divider(prs, "07", "AI COMPANION", "YOUR LOCALLY HOSTED BRAIN",
        "Three AI surfaces, one model, zero cloud calls. Everything you ask stays on your network."); n += 1
    slide_ai_companion(prs, n); n += 1

    # Notifications & Profile (combined slide)
    slide_notifications_profile(prs, n); n += 1

    # Tips
    slide_divider(prs, "09", "TIPS & BEST PRACTICE", "HABITS THAT COMPOUND",
        "Six habits that, over a quarter, make the workspace feel like a second brain."); n += 1
    slide_tips(prs, n); n += 1

    # Annex A
    slide_divider(prs, "A", "ANNEX · ARCHITECTURE", "HOW IT IS BUILT",
        "The local-first, three-tier system that puts you in full control of your data.", color=TEXT); n += 1
    slide_architecture(prs, n); n += 1

    # Annex B
    slide_divider(prs, "B", "ANNEX · PERMISSIONS", "WHO SEES WHAT",
        "The centralised role × group matrix that every screen in the workspace consults.", color=TEXT); n += 1
    slide_permission_matrix(prs, n); n += 1

    # Annex C
    slide_divider(prs, "C", "ANNEX · SHORTCUTS", "KEYBOARD + GLOSSARY",
        "For the keyboard-first user and anyone new to the DOINg.AI vocabulary.", color=TEXT); n += 1
    slide_shortcuts(prs, n); n += 1

    # Back cover
    slide_back_cover(prs)

    prs.save(output_path)
    import os
    size = os.path.getsize(output_path)
    print(f"OK — {output_path}")
    print(f"Slides: {len(prs.slides)}  ·  Size: {size/1024:.0f} KB")

if __name__ == "__main__":
    import pathlib
    out = str(pathlib.Path(__file__).resolve().parent.parent / "DOINg-AI-User-Guide.pptx")
    build(out)
